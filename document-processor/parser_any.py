# parser_any.py
#
# Версия 3.4: Универсальный роутер для всех типов файлов.
# --------------------------------------------------------------------------
import os
import sys
import json
import logging
import uuid
from typing import List, Dict, Any, Tuple, Callable

# --- Импорт всех парсеров ---
from parser_docx import parse_docx
from parser_txt import parse_txt
from parser_structured import parse_excel, parse_json, parse_xml

# --- Зависимости для встроенных парсеров ---
import pdfplumber
import pytesseract
from PIL import Image
from bs4 import BeautifulSoup
import pypandoc
from pptx import Presentation

# --- Вспомогательные функции ---
def make_error_block(doc_id: str, error_message: str) -> Tuple[List[Dict], Dict]:
    """Создает стандартный блок с ошибкой парсинга."""
    logging.error(f"[{doc_id}] {error_message}")
    blocks = [{"doc_id": doc_id, "chunk_id": 1, "type": "error", "text": error_message}]
    return blocks, {}

def get_filesystem_metadata(path: str) -> Dict[str, Any]:
    """Извлекает базовые метаданные из файловой системы."""
    import datetime
    try:
        return {
            "source_filename": os.path.basename(path),
            "created_fs": datetime.datetime.fromtimestamp(os.path.getctime(path)).isoformat(),
            "modified_fs": datetime.datetime.fromtimestamp(os.path.getmtime(path)).isoformat(),
            "size_bytes": os.path.getsize(path),
        }
    except Exception as e:
        logging.warning(f"Не удалось получить метаданные файловой системы для '{path}': {e}")
        return {"source_filename": os.path.basename(path)}

# --- Реализации встроенных парсеров ---

def parse_pdf(path: str, doc_id: str) -> Tuple[List[Dict], Dict]:
    blocks, chunk_id = [], 0
    doc_properties = {}
    try:
        with pdfplumber.open(path) as pdf:
            doc_properties = {**pdf.metadata, **get_filesystem_metadata(path)}
            for k, v in doc_properties.items():
                if isinstance(v, bytes):
                    try: doc_properties[k] = v.decode('utf-8', errors='ignore')
                    except: doc_properties[k] = str(v)

            for page_num, page in enumerate(pdf.pages, 1):
                text = page.extract_text(x_tolerance=2, y_tolerance=2) or ""
                if not text.strip() and os.getenv("OCR_ENABLED", "false").lower() == 'true':
                    try:
                        logging.info(f"PDF: [{doc_id}] Пустая страница {page_num}, применяю OCR...")
                        im = page.to_image(resolution=300).original
                        text = pytesseract.image_to_string(im, lang=os.getenv("OCR_LANG", "rus+eng"))
                    except Exception as ocr_e:
                        text = f"[OCR error on page {page_num}: {ocr_e}]"
                
                if text.strip():
                    chunk_id += 1
                    blocks.append({"doc_id": doc_id, "chunk_id": chunk_id, "section": f"Page {page_num}", "level": 0, "type": "paragraph", "text": text, "metadata": {"page": page_num}})
        return blocks, doc_properties
    except Exception as e:
        return make_error_block(doc_id, f"[PDF parse error: {e}]")

def parse_html(path: str, doc_id: str) -> Tuple[List[Dict], Dict]:
    blocks, chunk_id = [], 0
    doc_properties = get_filesystem_metadata(path)
    try:
        with open(path, "r", encoding="utf-8", errors='ignore') as f:
            content = f.read()
        
        soup = BeautifulSoup(content, "lxml")

        if soup.title and soup.title.string: doc_properties['title'] = soup.title.string.strip()
            
        main_content = soup.find('main') or soup.find('article') or soup.body or soup
        text = main_content.get_text(separator='\n\n', strip=True)
        for para in [p.strip() for p in text.split("\n\n") if p.strip()]:
            chunk_id += 1
            blocks.append({"doc_id": doc_id, "chunk_id": chunk_id, "text": para, "type": "paragraph", "section": doc_properties.get('title'), "level": 0, "metadata": {}})
        return blocks, doc_properties
    except Exception as e:
        return make_error_block(doc_id, f"[HTML parse error: {e}]")

def parse_pptx(path: str, doc_id: str) -> Tuple[List[Dict], Dict]:
    blocks, chunk_id = [], 0
    try:
        prs = Presentation(path)
        props = prs.core_properties
        doc_properties = {"author": props.author, "title": props.title, **get_filesystem_metadata(path)}
        for slide_idx, slide in enumerate(prs.slides, 1):
            slide_texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    slide_texts.append(shape.text_frame.text.strip())
            
            if slide_texts:
                chunk_id += 1
                full_slide_text = "\n\n".join(t for t in slide_texts if t)
                if full_slide_text:
                    blocks.append({"doc_id": doc_id, "chunk_id": chunk_id, "text": full_slide_text, "type": "slide_content", "section": f"Slide {slide_idx}", "level": 0, "metadata": {"slide": slide_idx}})
        return blocks, doc_properties
    except Exception as e:
        return make_error_block(doc_id, f"[PPTX parse error: {e}]")
        
def _wrap_parser(func: Callable, path: str, doc_id: str) -> Tuple[List[Dict], Dict]:
    """Вызывает парсер и унифицирует его вывод, добавляя метаданные файловой системы."""
    try:
        blocks, properties = func(path, doc_id)
        fs_meta = get_filesystem_metadata(path)
        final_properties = {**fs_meta, **(properties or {})}
        return blocks, final_properties
    except Exception as e:
        return make_error_block(doc_id, f"[Wrapped parser '{func.__name__}' error: {e}]")

def parse_any(path: str, doc_id: str) -> Tuple[List[Dict[str, Any]], Dict[str, Any]]:
    """
    Принимает путь к файлу и doc_id, вызывает нужный парсер
    и ВСЕГДА возвращает кортеж (блоки, свойства).
    """
    ext = os.path.splitext(path)[1].lower()
    logging.info(f"[{doc_id}] Выбор парсера для расширения '{ext}'...")
    
    PARSERS = {
        ".docx": parse_docx, ".pdf": parse_pdf, ".html": parse_html, ".htm": parse_html,
        ".pptx": parse_pptx, ".txt": parse_txt, ".xlsx": parse_excel, ".xls": parse_excel,
        ".json": parse_json, ".xml": parse_xml,
    }

    parser_func = PARSERS.get(ext)

    if parser_func:
        return _wrap_parser(parser_func, path, doc_id)
    else:
        logging.warning(f"[{doc_id}] Прямая поддержка для '{ext}' отсутствует. Используется fallback-парсер через Pandoc.")
        try:
            markdown_output = pypandoc.convert_file(path, 'md')
            doc_properties = get_filesystem_metadata(path)
            
            blocks = []
            if markdown_output.strip():
                blocks.append({"doc_id": doc_id, "chunk_id": 1, "text": markdown_output, "type": "paragraph", "section": None, "level": 0, "metadata": {}})
            return blocks, doc_properties
        except Exception as e:
             return make_error_block(doc_id, f"Unsupported file format: {ext}. Pandoc fallback also failed: {e}")

if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: python parser_any.py <file>")
        sys.exit(1)

    logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')
    test_filepath = sys.argv[1]
    test_doc_id = str(uuid.uuid4())
    
    try:
        blocks, properties = parse_any(test_filepath, test_doc_id)
        print("\n--- Document Properties ---")
        print(json.dumps(properties, ensure_ascii=False, indent=2))
        print(f"\n--- Parsed {len(blocks)} Blocks ---")
        for b in blocks:
            print(json.dumps(b, ensure_ascii=False, indent=2))
    except Exception as main_e:
        print(f"\n--- An error occurred during parsing: {main_e} ---")